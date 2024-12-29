#%%
import PyPDF2
from pptx import Presentation
from docx import Document
import io
from PIL import Image
from glob import glob
import os
from helper import *

class ExtractText:
    def __init__(self):
        self.system_prompt = "You are tasked with interpreting and describing images from parsed documents. When provided with an image and its surrounding context, generate a concise and accurate text interpretation of the image in no more than 50 tokens per image and maximum 300 tokens overall. Use the provided context"
        
    
    def extract_text(self, filepath):
        assert os.path.exists(filepath), "Filepath does not exist"
        if filepath.endswith(".pdf"): return self.read_pdf(filepath)
        elif filepath.endswith(".pptx"): return self.read_ppt(filepath)
        elif filepath.endswith(".txt"): return self.read_txt(filepath)
        elif filepath.endswith(".docx"): return self.read_docx(filepath)
        else: print("File format not supported")
        
    def read_pdf(self, pdf_path):
        with open(pdf_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            
            for page_number, page in enumerate(reader.pages, start=1):
                page_content = page.extract_text()
                images_base64 = []
                
                images = [page.images[0]] # Only consider the first image on the page since together api limits total images per url to 1
                for img_index, image in enumerate(images):
                    try:
                        im = Image.open(io.BytesIO(image.data))
                        buffer = io.BytesIO()
                        im.save(buffer, format="PNG") 
                        buffer.seek(0)
                        images_base64.append(base64.b64encode(buffer.getvalue()).decode("utf-8"))
                    except Exception as e:
                        #print(e)   
                        continue
                    
                if len(images_base64)>0:
                    retrieved_image_context = interpret_image_with_context(images_base64, self.system_prompt, page_content)
                    page_content = page_content + ". " +  retrieved_image_context 
                
                words = page_content.split()
                chunk_size = 350
                if len(words) <= 400:
                    chunk_size = 400
                
                chunks = [words[i:i + chunk_size] for i in range(0, len(words), chunk_size)]
                
                for chunk in chunks:
                    metadata = {
                        "content": ' '.join(chunk),
                        "page_number": page_number,
                        "file_type": "pdf",
                        "file_path": pdf_path
                    }
                    yield metadata
        
    def read_ppt(self, pptx_path):
        presentation = Presentation(pptx_path)
        
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            images_base64 = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text.append(paragraph.text)

                elif len(images_base64)==0 and shape.shape_type == 13:  # Only consider the first image on the page since together api limits total images per url to 1
                    try:
                        im = Image.open(io.BytesIO(shape.image.blob))
                        buffer = io.BytesIO()
                        im.save(buffer, format="PNG") 
                        buffer.seek(0)
                        images_base64.append(base64.b64encode(buffer.getvalue()).decode("utf-8"))
                    except Exception as e:
                        #print(e)
                        continue
                            
            if len(images_base64)>0:
                retrieved_image_context = interpret_image_with_context(images_base64, self.system_prompt, " ".join(slide_text))
                slide_text.append(retrieved_image_context)
                #print(retrieved_image_context)

            metadata = {
                "content": " ".join(slide_text),
                "page_number": slide_number,
                "file_type": "pptx",
                "file_path": pptx_path
            }
            yield metadata

    def read_txt(self, txt_path):
        with open(txt_path, 'r', encoding='utf-8') as txt_file:
            content = txt_file.read()
            words = content.split()
            chunk_size = 350
            if len(words) <= 400:
                chunk_size = 400
            
            chunks = [words[i:i + chunk_size] for i in range(0, len(words), chunk_size)]
            pages_metadata = []
            
            for chunk in chunks:
                metadata = {
                    "content": ' '.join(chunk),
                    "page_number": 1,
                    "file_type": "txt",
                    "file_path": txt_path
                }
                yield metadata
    
    def read_docx(self, docx_path):
        doc = Document(docx_path)
        content = []

        for element in doc.element.body:
            if element.tag.endswith('p'):
                para = element.text
                if para.strip():
                    content.append(para)
            elif element.tag.endswith('tbl'):
                # TODO: Handle tables 
                pass
            elif element.tag.endswith('drawing'):
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        pass #TODO: Handle images

        full_text = ' '.join(content)
        words = full_text.split()
        chunk_size = 350
        if len(words) <= 400:
            chunk_size = 400
        
        chunks = [words[i:i + chunk_size] for i in range(0, len(words), chunk_size)]
        
        page_number = 1
        for chunk in chunks:
            metadata = {
                "content": ' '.join(chunk),
                "page_number": page_number,
                "file_type": "docx",
                "file_path": docx_path
            }
            yield metadata
            page_number += 1
